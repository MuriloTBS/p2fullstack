from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

LOGO_PATH = "/Users/sergionogueira/Desktop/p2fullstack/logo s4 -02png.png"
HAS_LOGO  = os.path.exists(LOGO_PATH)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Paleta ──────────────────────────────────────────────
INDIGO      = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_DARK = RGBColor(0x31, 0x27, 0xAE)
INDIGO_BG   = RGBColor(0x1E, 0x1B, 0x7A)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_50     = RGBColor(0xF9, 0xFA, 0xFB)
GRAY_100    = RGBColor(0xF3, 0xF4, 0xF6)
GRAY_200    = RGBColor(0xE5, 0xE7, 0xEB)
GRAY_400    = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_600    = RGBColor(0x4B, 0x55, 0x63)
GRAY_700    = RGBColor(0x37, 0x41, 0x51)
GREEN       = RGBColor(0x05, 0x96, 0x69)
GREEN_LIGHT = RGBColor(0xD1, 0xFA, 0xE5)
YELLOW      = RGBColor(0xD9, 0x77, 0x06)
YELLOW_L    = RGBColor(0xFE, 0xF3, 0xC7)
PURPLE      = RGBColor(0x6D, 0x28, 0xD9)
PURPLE_L    = RGBColor(0xED, 0xE9, 0xFE)
BLUE        = RGBColor(0x02, 0x84, 0xC7)
BLUE_L      = RGBColor(0xE0, 0xF2, 0xFE)
RED         = RGBColor(0xDC, 0x26, 0x26)

BLANK = prs.slide_layouts[6]

# ── helpers ─────────────────────────────────────────────
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill, line=None, lw=Pt(1)):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = lw
    else:
        sh.line.fill.background()
    return sh

def label(slide, text, l, t, w, h, size=14, bold=False,
          color=WHITE, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def hline(slide, t, color=INDIGO, l=0.65, w=12.0, h=0.045):
    box(slide, l, t, w, h, color)

def header(slide, title, subtitle="", bar_color=INDIGO_DARK):
    box(slide, 0, 0, 13.33, 1.45, bar_color)
    label(slide, title,    0.6, 0.2,  12.0, 0.7,  size=30, bold=True,  color=WHITE)
    if subtitle:
        label(slide, subtitle, 0.6, 0.9,  12.0, 0.4,  size=13, color=RGBColor(0xC7,0xD2,0xFE))
    hline(slide, 1.6)

def tag(slide, text, l, t, w=1.4, h=0.3, fill=INDIGO, tc=WHITE, size=10):
    box(slide, l, t, w, h, fill)
    label(slide, text, l, t, w, h, size=size, bold=True, color=tc, align=PP_ALIGN.CENTER)

def check_list(slide, items, l, t, size=12, color=GRAY_700, gap=0.38):
    y = t
    for item in items:
        label(slide, f"  +  {item}", l, y, 12.0, gap, size=size, color=color)
        y += gap

# ════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, INDIGO_DARK)

# barra superior
box(s, 0, 0, 13.33, 1.6, RGBColor(0x1E,0x1B,0x7A))
# acento lateral esquerdo
box(s, 0, 0, 0.18, 7.5, INDIGO)

# logo quadrado
box(s, 0.65, 0.28, 1.05, 1.05, INDIGO)
label(s, "P2", 0.65, 0.37, 1.05, 0.65, size=26, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)

# cabeçalho
label(s, "UNIVASSOURAS  —  MARICA", 1.9, 0.3,  10, 0.4,
      size=11, color=RGBColor(0xA5,0xB4,0xFC))
label(s, "Laboratorio de Programacao Full Stack", 1.9, 0.7, 10, 0.4,
      size=14, color=RGBColor(0xC7,0xD2,0xFE), bold=True)

hline(s, 1.72, RGBColor(0x6D,0x66,0xEA), l=0.65, w=11.9)

# titulo
label(s, "Atividade P2  —  Full Stack", 0.65, 1.9, 12, 0.5,
      size=15, color=RGBColor(0xA5,0xB4,0xFC))
label(s, "Sistema de Gestao com\nReact + Laravel + Docker", 0.65, 2.38, 12, 1.8,
      size=40, bold=True, color=WHITE)

# tags tecnologia
tx = 0.65
for t_lbl, t_col in [("React 18", INDIGO), ("Tailwind CSS", BLUE),
                      ("Laravel 12", GREEN), ("Docker", PURPLE)]:
    tag(s, t_lbl, tx, 4.5, w=1.7, h=0.38, fill=t_col, size=12)
    tx += 1.85

# rodape
box(s, 0, 6.3, 13.33, 1.2, INDIGO_BG)
label(s, "Aluno:     Sergio Murilo",       0.8, 6.4,  6, 0.38, size=14, bold=True, color=WHITE)
label(s, "Professor: Fabricio Dias",       0.8, 6.78, 6, 0.35, size=13, color=RGBColor(0xC7,0xD2,0xFE))
label(s, "2026", 11.0, 6.45, 2.2, 0.55, size=28, bold=True,
      color=RGBColor(0x6D,0x66,0xEA), align=PP_ALIGN.RIGHT)
# logo S4 — canto superior direito da capa
if HAS_LOGO:
    s.shapes.add_picture(LOGO_PATH, Inches(11.2), Inches(0.25), height=Inches(0.95))


# ════════════════════════════════════════════════════════
# SLIDE 2 — OBJETIVO
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_50)
header(s, "Objetivo do Projeto",
       "Criar uma aplicacao fullstack completa, integrada e containerizada com Docker")

label(s, "Desenvolver uma aplicacao web moderna que integre frontend em React com "
         "backend Laravel via API REST, com autenticacao, controle de acesso por perfil e "
         "execucao completa via Docker Compose, sem dependencia de ambiente local.",
      0.6, 1.75, 12.1, 0.9, size=14, color=GRAY_700, wrap=True)

cards = [
    ("Autenticacao",      "Login seguro com\ntoken JWT Bearer",          INDIGO,  RGBColor(0xEE,0xF2,0xFF)),
    ("CRUD Completo",     "Categorias e Usuarios\ncom regras de acesso",  GREEN,   GREEN_LIGHT),
    ("Containerizacao",   "Toda a stack Docker\nsem setup manual",        BLUE,    BLUE_L),
    ("Interface Moderna", "React + Tailwind CSS\nresponsivo e limpo",     PURPLE,  PURPLE_L),
]
cx = 0.5
for title, desc, border, bg2 in cards:
    box(s, cx, 2.9, 2.9, 2.85, bg2, border, lw=Pt(2))
    box(s, cx, 2.9, 2.9, 0.42, border)
    label(s, title, cx+0.12, 2.94, 2.65, 0.35,
          size=13, bold=True, color=WHITE)
    label(s, desc, cx+0.15, 3.45, 2.65, 0.9,
          size=13, color=GRAY_700, wrap=True)
    cx += 3.08


# ════════════════════════════════════════════════════════
# SLIDE 3 — TECNOLOGIAS
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_50)
header(s, "Tecnologias Utilizadas", "Stack moderno e profissional em todas as camadas")

cols = [
    ("FRONTEND",        INDIGO, [
        ("React 18",           "Biblioteca de UI com componentes reutilizaveis"),
        ("Tailwind CSS",       "Estilizacao utility-first e responsiva"),
        ("Vite",               "Build tool ultrarapido para desenvolvimento"),
        ("Axios",              "Cliente HTTP com interceptors para JWT"),
        ("React Router v6",    "Navegacao SPA com rotas protegidas"),
    ]),
    ("BACKEND",         GREEN, [
        ("PHP 8.3",            "Linguagem do servidor"),
        ("Laravel 12",         "Framework MVC com API REST nativo"),
        ("MySQL 8.0",          "Banco de dados relacional"),
        ("Token Auth",         "Guard nativo do Laravel (sem pacote extra)"),
        ("Apache",             "Servidor web no container"),
    ]),
    ("INFRAESTRUTURA",  BLUE, [
        ("Docker",             "Containerizacao de todos os servicos"),
        ("Docker Compose",     "Orquestracao multi-container"),
        ("nginx",              "Serve o React + proxy /api/ -> backend"),
        ("Composer",           "Gerenciador de pacotes PHP"),
        ("npm",                "Gerenciador de pacotes JavaScript"),
    ]),
]

cx = 0.38
for title, color, items in cols:
    box(s, cx, 1.75, 4.05, 0.44, color)
    label(s, title, cx, 1.79, 4.05, 0.36,
          size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    ry = 2.26
    for name, desc in items:
        box(s, cx, ry, 4.05, 0.78, WHITE, color, lw=Pt(1))
        label(s, name, cx+0.14, ry+0.06, 3.8, 0.32, size=12, bold=True, color=INDIGO_DARK)
        label(s, desc, cx+0.14, ry+0.38, 3.8, 0.34, size=10, color=GRAY_400)
        ry += 0.84
    cx += 4.32


# ════════════════════════════════════════════════════════
# SLIDE 4 — ARQUITETURA
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_50)
header(s, "Arquitetura do Sistema",
       "Tres containers Docker se comunicam via rede interna")

# --- caixas do diagrama ---
def arc_box(sl, l, t, w, h, title, sub, fill, tc=WHITE):
    box(sl, l, t, w, h, fill)
    label(sl, title, l+0.12, t+0.1,  w-0.2, 0.38, size=13, bold=True, color=tc)
    label(sl, sub,   l+0.12, t+0.48, w-0.2, 0.36, size=10,
          color=RGBColor(0xE0,0xE7,0xFF) if tc==WHITE else GRAY_400, italic=True)

arc_box(s, 0.4,  2.55, 2.3, 1.0, "NAVEGADOR",        "http://localhost:3000", BLUE)
arc_box(s, 3.7,  1.55, 3.0, 1.0, "React Build",       "arquivos estaticos",    GRAY_400)
arc_box(s, 3.7,  2.9,  3.0, 1.0, "nginx  (frontend)", "porta 3000",            INDIGO)
arc_box(s, 3.7,  4.55, 3.0, 1.0, "Laravel  (backend)","porta 8080",            GREEN)
arc_box(s, 8.2,  4.55, 2.8, 1.0, "MySQL 8.0",         "rede interna Docker",   YELLOW,
        tc=GRAY_700)

# setas texto
label(s, ">>>", 2.78, 2.9,  0.85, 0.42, size=14, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
label(s, ":3000",2.6,  2.62, 1.1,  0.28, size=9,  color=GRAY_400, align=PP_ALIGN.CENTER)

label(s, "v", 4.98, 1.65, 0.6, 0.28, size=14, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
label(s, "serve estaticos", 3.7, 1.62, 3.0, 0.26, size=8, color=GRAY_400, align=PP_ALIGN.CENTER)

label(s, "v", 4.98, 4.05, 0.6, 0.35, size=14, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
label(s, "proxy  /api/*", 3.7, 4.05, 3.0, 0.3, size=9, color=GRAY_400, align=PP_ALIGN.CENTER)

label(s, ">>>", 6.8, 4.88, 1.25, 0.38, size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
label(s, "SQL query", 6.8, 4.62, 1.25, 0.26, size=8, color=GRAY_400, align=PP_ALIGN.CENTER)

# painel lateral
box(s, 11.25, 1.62, 1.9, 4.2, WHITE, INDIGO, lw=Pt(2))
label(s, "Rede Interna\nDocker", 11.25, 1.7, 1.9, 0.55,
      size=11, bold=True, color=INDIGO_DARK, align=PP_ALIGN.CENTER)
hline(s, 2.32, INDIGO, l=11.3, w=1.8, h=0.03)
notas = ["Containers isolados","","Comunicacao via","nome do servico","","MySQL nao exposto","ao host","","nginx proxy","transparente"]
ny = 2.4
for n in notas:
    label(s, n, 11.35, ny, 1.75, 0.28, size=9, color=GRAY_700)
    ny += 0.28


# ════════════════════════════════════════════════════════
# SLIDE 5 — BACKEND
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_50)
header(s, "Backend — Laravel 12 (API REST)",
       "API RESTful com autenticacao por token e controle de acesso por perfil",
       bar_color=RGBColor(0x04,0x5C,0x42))
hline(s, 1.6, GREEN)

# --- Auth ---
box(s, 0.5, 1.75, 5.7, 2.25, WHITE, GREEN, lw=Pt(2))
box(s, 0.5, 1.75, 5.7, 0.44, GREEN)
label(s, "AUTENTICACAO", 0.62, 1.8, 5.5, 0.35, size=12, bold=True, color=WHITE)
auth_items = [
    "Guard token nativo do Laravel (sem Sanctum)",
    "Token Bearer de 60 caracteres aleatorios",
    "Armazenado na coluna api_token dos usuarios",
    "Invalidado no banco de dados no logout",
    "Protecao via middleware  auth:api",
]
ay = 2.26
for it in auth_items:
    label(s, f"   +  {it}", 0.62, ay, 5.4, 0.3, size=11, color=GRAY_700)
    ay += 0.31

# --- Controle de acesso ---
box(s, 0.5, 4.12, 5.7, 1.58, WHITE, GREEN, lw=Pt(2))
box(s, 0.5, 4.12, 5.7, 0.44, GREEN)
label(s, "CONTROLE DE ACESSO", 0.62, 4.17, 5.5, 0.35, size=12, bold=True, color=WHITE)
roles_data = [
    ("ADMIN",   "Acesso total — CRUD Usuarios + Categorias", YELLOW,   GRAY_700),
    ("USER",    "Acesso parcial — apenas CRUD de Categorias", BLUE,     WHITE),
]
ry2 = 4.65
for role, desc, rc, rtc in roles_data:
    box(s, 0.62, ry2, 0.9, 0.34, rc)
    label(s, role, 0.62, ry2, 0.9, 0.34, size=10, bold=True, color=rtc, align=PP_ALIGN.CENTER)
    label(s, desc, 1.6,  ry2, 4.5, 0.34, size=11, color=GRAY_700)
    ry2 += 0.44

# --- Rotas ---
box(s, 6.5, 1.75, 6.3, 5.45, WHITE, GREEN, lw=Pt(2))
box(s, 6.5, 1.75, 6.3, 0.44, GREEN)
label(s, "ROTAS DA API", 6.62, 1.8, 6.1, 0.35, size=12, bold=True, color=WHITE)

routes = [
    ("POST",   "/api/login",            "Publico", YELLOW,   GRAY_700),
    ("POST",   "/api/logout",           "Auth",    GREEN,    WHITE),
    ("GET",    "/api/me",               "Auth",    GREEN,    WHITE),
    ("GET",    "/api/categorias",       "Auth",    GREEN,    WHITE),
    ("POST",   "/api/categorias",       "Auth",    GREEN,    WHITE),
    ("PUT",    "/api/categorias/{id}",  "Auth",    GREEN,    WHITE),
    ("DELETE", "/api/categorias/{id}",  "Auth",    GREEN,    WHITE),
    ("GET",    "/api/users",            "Admin",   INDIGO,   WHITE),
    ("POST",   "/api/users",            "Admin",   INDIGO,   WHITE),
    ("PUT",    "/api/users/{id}",       "Admin",   INDIGO,   WHITE),
    ("DELETE", "/api/users/{id}",       "Admin",   INDIGO,   WHITE),
]
mc = {"GET": GREEN, "POST": INDIGO, "PUT": YELLOW, "DELETE": RED}
rry = 2.26
for method, path, guard, gc, gtc in routes:
    box(s, 6.62, rry, 0.78, 0.3, mc.get(method, INDIGO))
    label(s, method, 6.62, rry, 0.78, 0.3, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    label(s, path,   7.46, rry, 3.6,  0.3, size=10, color=GRAY_700)
    box(s, 11.12, rry, 1.0, 0.3, gc)
    label(s, guard, 11.12, rry, 1.0, 0.3, size=9, bold=True, color=gtc, align=PP_ALIGN.CENTER)
    rry += 0.35


# ════════════════════════════════════════════════════════
# SLIDE 6 — FRONTEND
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_50)
header(s, "Frontend — React 18 + Tailwind CSS",
       "SPA moderna com gerenciamento de estado, rotas protegidas e consumo de API")

# --- estrutura ---
box(s, 0.5, 1.75, 3.9, 5.3, WHITE, INDIGO, lw=Pt(2))
box(s, 0.5, 1.75, 3.9, 0.44, INDIGO)
label(s, "ESTRUTURA DE PASTAS", 0.62, 1.8, 3.7, 0.35, size=11, bold=True, color=WHITE)
tree = [
    ("src/",                 True),
    ("   components/",       True),
    ("      Alert.jsx",      False),
    ("      Layout.jsx",     False),
    ("      Navbar.jsx",     False),
    ("      PrivateRoute.jsx",False),
    ("   context/",          True),
    ("      AuthContext.jsx", False),
    ("   services/",         True),
    ("      api.js",         False),
    ("   pages/",            True),
    ("      Login.jsx",      False),
    ("      Dashboard.jsx",  False),
    ("      categorias/",    True),
    ("      usuarios/",      True),
    ("   App.jsx",           False),
]
ty2 = 2.26
for lbl2, is_dir in tree:
    c2 = INDIGO_DARK if is_dir else GRAY_700
    label(s, lbl2, 0.62, ty2, 3.65, 0.27, size=10, bold=is_dir, color=c2)
    ty2 += 0.27

# --- cards componentes ---
comps = [
    ("AuthContext",   "Estado global do usuario logado.\nFuncoes login(), logout(), isAdmin()",  INDIGO,  RGBColor(0xEE,0xF2,0xFF)),
    ("api.js",        "Axios com interceptors.\nBearer token automatico em todo request",        GREEN,   GREEN_LIGHT),
    ("PrivateRoute",  "Bloqueia rotas sem autenticacao.\nAdminRoute barra usuarios comuns",      PURPLE,  PURPLE_L),
    ("Navbar",        "Menu responsivo com badge ADMIN\ne botao de Logout",                      BLUE,    BLUE_L),
]
cx2 = 4.55
cy2 = 1.75
for i, (title, desc, border, bgc) in enumerate(comps):
    if i == 2:
        cx2 = 4.55
        cy2 = 4.0
    box(s, cx2, cy2, 4.25, 1.95, bgc, border, lw=Pt(2))
    box(s, cx2, cy2, 4.25, 0.44, border)
    label(s, title, cx2+0.12, cy2+0.06, 4.0, 0.35, size=13, bold=True, color=WHITE)
    label(s, desc,  cx2+0.15, cy2+0.6,  4.0, 1.2,  size=12, color=GRAY_700)
    cx2 += 4.42

# --- rotas SPA ---
box(s, 4.55, 6.05, 8.55, 0.65, WHITE, INDIGO, lw=Pt(1))
label(s, "Rotas SPA:", 4.68, 6.13, 2.0, 0.3, size=11, bold=True, color=INDIGO_DARK)
label(s, "/login   /dashboard   /categorias   /categorias/novo   /usuarios   /usuarios/novo",
      6.7, 6.14, 6.3, 0.3, size=10, color=GRAY_700)


# ════════════════════════════════════════════════════════
# SLIDE 7 — DOCKER
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_50)
header(s, "Docker & Docker Compose",
       "Toda a stack containerizada — sem necessidade de ambiente local configurado",
       bar_color=BLUE)
hline(s, 1.6, BLUE)

cont_data = [
    ("mysql", "MySQL 8.0",
     ["Imagem: mysql:8.0",
      "Volume persistente: mysql_data",
      "Healthcheck: mysqladmin ping",
      "Backend aguarda ready",
      "Nao exposto ao host",
      "Variaveis de ambiente:","  DB_DATABASE, USER, PASS"],
     YELLOW, GRAY_700),
    ("backend","Laravel + Apache",
     ["Imagem: php:8.3-apache",
      "Dockerfile: ./backend/",
      "Depende de: mysql (healthy)",
      "entrypoint.sh faz:",
      "  1. Aguarda MySQL",
      "  2. php artisan migrate",
      "  3. php artisan db:seed",
      "  4. apache2-foreground"],
     GREEN, WHITE),
    ("frontend","React + nginx",
     ["Dockerfile multi-stage:",
      "  Stage 1: node:20-alpine",
      "    npm install + npm run build",
      "  Stage 2: nginx:alpine",
      "    serve /dist estatico",
      "nginx.conf:",
      "  proxy /api/* -> backend",
      "  SPA fallback index.html"],
     INDIGO, WHITE),
]
cx3 = 0.38
for svc, subtitle, lines, color, tc in cont_data:
    box(s, cx3, 1.75, 4.1, 4.85, WHITE, color, lw=Pt(2))
    box(s, cx3, 1.75, 4.1, 0.82, color)
    label(s, svc,      cx3+0.12, 1.82, 3.85, 0.38, size=16, bold=True, color=WHITE)
    label(s, subtitle, cx3+0.12, 2.2,  3.85, 0.32, size=11, color=RGBColor(0xE0,0xE7,0xFF), italic=True)
    ly3 = 2.68
    for ln in lines:
        is_key = not ln.startswith("  ")
        label(s, ln, cx3+0.15, ly3, 3.8, 0.3, size=10,
              color=color if is_key else GRAY_700, bold=is_key)
        ly3 += 0.3
    cx3 += 4.32

# portas
box(s, 0.38, 6.75, 12.57, 0.55, WHITE, BLUE, lw=Pt(1))
label(s, "Portas:", 0.55, 6.82, 1.3, 0.3, size=12, bold=True, color=BLUE)
ports_info = [
    ("frontend  :3000", "http://localhost:3000", INDIGO),
    ("backend   :8080", "http://localhost:8080",  GREEN),
    ("mysql     :---",  "somente rede interna",    YELLOW),
]
px2 = 1.95
for p_lbl, p_url, pc in ports_info:
    label(s, p_lbl, px2, 6.8,  2.8, 0.3, size=11, bold=True, color=pc)
    label(s, p_url, px2, 7.08, 2.8, 0.28, size=9, color=GRAY_400)
    px2 += 3.9


# ════════════════════════════════════════════════════════
# SLIDE 8 — FUNCIONALIDADES
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_50)
header(s, "Funcionalidades — Telas do Sistema",
       "Interface completa com feedback visual em todas as operacoes")

screens_data = [
    ("Login",       "Autenticacao JWT",
     ["E-mail e senha validados","Token salvo no localStorage",
      "Mensagem de erro clara","Credenciais de teste visiveis"],
     INDIGO, RGBColor(0xEE,0xF2,0xFF)),
    ("Dashboard",   "Visao geral",
     ["Totais de categorias/usuarios","Cards clicaveis","Saudacao com nome","Badge ADMIN"],
     GREEN,  GREEN_LIGHT),
    ("Categorias",  "CRUD completo",
     ["Listar em tabela responsiva","Cadastrar nova categoria","Editar com formulario","Confirmar exclusao"],
     PURPLE, PURPLE_L),
    ("Usuarios",    "Somente admin",
     ["Badge Admin / Usuario","Bloqueia nao-admins","Nao remove conta propria","Senha opcional na edicao"],
     YELLOW, YELLOW_L),
]
cx4 = 0.4
for title, sub, feats, color, bgc in screens_data:
    box(s, cx4, 1.75, 2.98, 5.35, bgc, color, lw=Pt(2))
    box(s, cx4, 1.75, 2.98, 0.78, color)
    label(s, title, cx4+0.12, 1.82, 2.75, 0.38, size=15, bold=True, color=WHITE)
    label(s, sub,   cx4+0.12, 2.2,  2.75, 0.28, size=10,
          color=RGBColor(0xE0,0xE7,0xFF) if color!=YELLOW else GRAY_700, italic=True)
    fy2 = 2.68
    for feat in feats:
        box(s, cx4+0.18, fy2+0.06, 0.22, 0.22, color)
        label(s, feat, cx4+0.5, fy2, 2.35, 0.36, size=11.5, color=GRAY_700)
        fy2 += 0.44
    cx4 += 3.18


# ════════════════════════════════════════════════════════
# SLIDE 9 — COMO EXECUTAR
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GRAY_50)
header(s, "Como Executar o Projeto",
       "Apenas Docker necessario — sem Node.js, PHP ou Composer no host")

steps_data = [
    ("1", "Clone o repositorio",
     "git clone https://github.com/MuriloTBS/p2fullstack.git    cd p2fullstack",
     INDIGO),
    ("2", "Suba todos os containers",
     "docker compose up --build",
     GREEN),
    ("3", "Aguarde a inicializacao automatica",
     "MySQL inicia  ->  migrations  ->  seeders  ->  Apache + nginx prontos",
     BLUE),
    ("4", "Acesse no navegador",
     "http://localhost:3000          Login: admin@example.com  /  password",
     PURPLE),
]
sy2 = 1.78
for num, title, cmd, color in steps_data:
    box(s, 0.5,  sy2, 0.72, 0.72, color)
    label(s, num, 0.5, sy2, 0.72, 0.72, size=24, bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 1.35, sy2, 11.45, 0.72, WHITE, color, lw=Pt(1))
    label(s, title, 1.5,  sy2+0.06, 4.5, 0.3, size=13, bold=True, color=INDIGO_DARK)
    label(s, cmd,   1.5,  sy2+0.37, 11.1,0.3, size=11, color=GRAY_700, italic=True)
    sy2 += 0.9

box(s, 0.5, 5.65, 12.3, 1.2, WHITE, INDIGO, lw=Pt(2))
label(s, "Pre-requisitos:", 0.65, 5.73, 3.5, 0.35, size=13, bold=True, color=INDIGO_DARK)
label(s, "Docker >= 24   |   Docker Compose v2",
      4.3, 5.74, 8.4, 0.35, size=13, color=GRAY_600)
label(s, "Nao e necessario instalar Node.js, PHP, Composer ou MySQL na maquina — "
         "tudo e gerenciado pelos containers.",
      0.65, 6.12, 12.0, 0.6, size=11, color=GRAY_400)


# ════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSAO
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, INDIGO_DARK)
box(s, 0, 0, 13.33, 7.5, RGBColor(0x31,0x27,0xAE))
box(s, 0, 5.55, 13.33, 1.95, INDIGO_BG)
box(s, 0, 0, 0.22, 7.5, INDIGO)

label(s, "CONCLUSAO", 0.65, 0.42, 12, 0.42, size=13, bold=True,
      color=RGBColor(0xA5,0xB4,0xFC))
label(s, "Todos os objetivos\nalcancados", 0.65, 0.82, 12, 1.3,
      size=42, bold=True, color=WHITE)
hline(s, 2.32, RGBColor(0x6D,0x66,0xEA), l=0.65, w=11.9)

conquistas = [
    "Frontend React com Tailwind CSS — interface moderna e responsiva",
    "Backend Laravel 12 com API REST e autenticacao por token Bearer",
    "CRUD completo de Categorias e Usuarios com regras de acesso",
    "Controle de acesso por perfil: admin acessa tudo, user acessa categorias",
    "Containerizacao completa: MySQL + Laravel + React via Docker Compose",
    "nginx como proxy transparente — zero configuracao adicional no host",
    "Migrations e seeders automaticos no startup dos containers",
    "Codigo e documentacao publicados no GitHub",
]
cy3 = 2.52
for c in conquistas:
    box(s, 0.65, cy3+0.06, 0.26, 0.26, INDIGO)
    label(s, c, 1.1, cy3, 11.5, 0.36, size=13, color=WHITE)
    cy3 += 0.37

# rodape
label(s, "github.com/MuriloTBS/p2fullstack", 0.65, 5.68, 9.0, 0.45,
      size=14, bold=True, color=RGBColor(0xA5,0xB4,0xFC))
label(s, "Laboratorio de Programacao Full Stack   |   Prof. Fabricio Dias   |   Univassouras — Marica",
      0.65, 6.14, 12.0, 0.35, size=11, color=RGBColor(0xC7,0xD2,0xFE))
label(s, "Aluno: Sergio Murilo   |   2026",
      0.65, 6.55, 12.0, 0.38, size=13, bold=True, color=WHITE)
# logo S4 — canto inferior direito da conclusao
if HAS_LOGO:
    s.shapes.add_picture(LOGO_PATH, Inches(11.0), Inches(6.2), height=Inches(1.0))


# ── salva ────────────────────────────────────────────────
out = "/Users/sergionogueira/Desktop/p2fullstack/Apresentacao_P2_FullStack.pptx"
prs.save(out)
print(f"Salvo: {out}")
